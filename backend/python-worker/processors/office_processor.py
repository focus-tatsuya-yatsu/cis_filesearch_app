"""
Office Processor Module
Handles Microsoft Office documents (Word, Excel, PowerPoint)
"""

import io
import time
from pathlib import Path
from typing import Optional

import docx
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image

from .base_processor import BaseProcessor, ProcessingResult


class OfficeProcessor(BaseProcessor):
    """
    Processor for Microsoft Office documents
    Supports: DOC, DOCX, XLS, XLSX, PPT, PPTX
    """

    def can_process(self, file_path: str) -> bool:
        """Check if this processor can handle the file"""
        ext = Path(file_path).suffix.lower()
        return ext in self.config.file_types.office_extensions

    def process(self, file_path: str) -> ProcessingResult:
        """
        Process Office document

        Args:
            file_path: Path to Office document

        Returns:
            ProcessingResult object
        """
        start_time = time.time()

        try:
            # Validate file size
            if not self._validate_file_size(
                file_path,
                self.config.processing.max_office_size_mb
            ):
                return self._create_error_result(
                    file_path,
                    f"File exceeds maximum size of {self.config.processing.max_office_size_mb}MB"
                )

            # Get file info
            file_info = self._get_file_info(file_path)
            file_ext = file_info['file_type']

            # Route to appropriate handler
            if file_ext in {'.docx', '.doc'}:
                extracted_text, metadata = self._process_word(file_path)
            elif file_ext in {'.xlsx', '.xls'}:
                extracted_text, metadata = self._process_excel(file_path)
            elif file_ext in {'.pptx', '.ppt'}:
                extracted_text, metadata = self._process_powerpoint(file_path)
            else:
                return self._create_error_result(
                    file_path,
                    f"Unsupported Office format: {file_ext}"
                )

            # Add file metadata
            metadata.update(self._extract_metadata(file_path))

            # Generate thumbnail if enabled
            thumbnail_data = None
            if self.config.thumbnail.generate_for_office:
                thumbnail_data = self._generate_thumbnail(file_path, file_ext)

            # Calculate processing time
            processing_time = time.time() - start_time

            # Create result
            result = ProcessingResult(
                success=True,
                file_path=file_info['file_path'],
                file_name=file_info['file_name'],
                file_size=file_info['file_size'],
                file_type=file_info['file_type'],
                mime_type=self.config.file_types.get_mime_type(file_info['file_type']),
                extracted_text=extracted_text,
                word_count=self._count_words(extracted_text),
                char_count=len(extracted_text),
                page_count=metadata.get('page_count', 0),
                thumbnail_data=thumbnail_data,
                thumbnail_format=self.config.thumbnail.thumbnail_format,
                metadata=metadata,
                processing_time_seconds=processing_time,
                processor_name=self.__class__.__name__,
            )

            self.logger.info(
                f"Successfully processed Office document: {file_info['file_name']} "
                f"({len(extracted_text)} chars in {processing_time:.2f}s)"
            )

            return result

        except Exception as e:
            self.logger.error(f"Error processing Office document: {e}", exc_info=True)
            return self._create_error_result(
                file_path,
                f"Processing error: {str(e)}"
            )

    def _process_word(self, file_path: str) -> tuple[str, dict]:
        """
        Process Word document

        Args:
            file_path: Path to Word document

        Returns:
            Tuple of (extracted_text, metadata)
        """
        try:
            doc = docx.Document(file_path)

            # Extract text from paragraphs
            text_parts = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_parts.append(paragraph.text)

            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text)
                    if row_text:
                        text_parts.append('\t'.join(row_text))

            extracted_text = '\n'.join(text_parts)

            # Extract metadata
            metadata = {
                'document_type': 'word',
                'paragraph_count': len(doc.paragraphs),
                'table_count': len(doc.tables),
                'page_count': 0,  # Word doesn't provide page count easily
            }

            # Try to get core properties
            try:
                core_props = doc.core_properties
                metadata['title'] = core_props.title or ''
                metadata['author'] = core_props.author or ''
                metadata['subject'] = core_props.subject or ''
                metadata['keywords'] = core_props.keywords or ''
                metadata['last_modified_by'] = core_props.last_modified_by or ''
            except:
                pass

            return extracted_text, metadata

        except Exception as e:
            self.logger.error(f"Failed to process Word document: {e}")
            raise

    def _process_excel(self, file_path: str) -> tuple[str, dict]:
        """
        Process Excel spreadsheet

        Args:
            file_path: Path to Excel file

        Returns:
            Tuple of (extracted_text, metadata)
        """
        try:
            workbook = load_workbook(file_path, read_only=True, data_only=True)

            text_parts = []
            total_rows = 0
            total_cells = 0

            # Process each sheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]

                # Add sheet name as header
                text_parts.append(f"=== Sheet: {sheet_name} ===")

                # Extract cell values
                for row in sheet.iter_rows(values_only=True):
                    total_rows += 1
                    row_text = []

                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value))
                            total_cells += 1

                    if row_text:
                        text_parts.append('\t'.join(row_text))

            extracted_text = '\n'.join(text_parts)

            # Metadata
            metadata = {
                'document_type': 'excel',
                'sheet_count': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames,
                'row_count': total_rows,
                'cell_count': total_cells,
            }

            # Try to get workbook properties
            try:
                props = workbook.properties
                metadata['title'] = props.title or ''
                metadata['author'] = props.creator or ''
                metadata['subject'] = props.subject or ''
                metadata['keywords'] = props.keywords or ''
            except:
                pass

            workbook.close()

            return extracted_text, metadata

        except Exception as e:
            self.logger.error(f"Failed to process Excel file: {e}")
            raise

    def _process_powerpoint(self, file_path: str) -> tuple[str, dict]:
        """
        Process PowerPoint presentation

        Args:
            file_path: Path to PowerPoint file

        Returns:
            Tuple of (extracted_text, metadata)
        """
        try:
            presentation = Presentation(file_path)

            text_parts = []
            slide_count = len(presentation.slides)

            # Process each slide
            for i, slide in enumerate(presentation.slides, 1):
                # Add slide header
                text_parts.append(f"=== Slide {i} ===")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                    # Extract text from tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text)
                            if row_text:
                                text_parts.append('\t'.join(row_text))

            extracted_text = '\n'.join(text_parts)

            # Metadata
            metadata = {
                'document_type': 'powerpoint',
                'slide_count': slide_count,
                'page_count': slide_count,
            }

            # Try to get core properties
            try:
                core_props = presentation.core_properties
                metadata['title'] = core_props.title or ''
                metadata['author'] = core_props.author or ''
                metadata['subject'] = core_props.subject or ''
                metadata['keywords'] = core_props.keywords or ''
            except:
                pass

            return extracted_text, metadata

        except Exception as e:
            self.logger.error(f"Failed to process PowerPoint file: {e}")
            raise

    def _generate_thumbnail(self, file_path: str, file_ext: str) -> Optional[bytes]:
        """
        Generate thumbnail for Office document

        For PPTX: Extract first slide image
        For DOCX/XLSX: Convert first page to image using LibreOffice (if available)

        Args:
            file_path: Path to Office document
            file_ext: File extension

        Returns:
            Thumbnail as bytes or None
        """
        try:
            # PowerPoint: Extract thumbnail from presentation
            if file_ext in {'.pptx', '.ppt'}:
                return self._generate_pptx_thumbnail(file_path)

            # Word/Excel: Try LibreOffice conversion
            else:
                return self._generate_office_thumbnail_via_libreoffice(file_path)

        except Exception as e:
            self.logger.warning(f"Thumbnail generation failed for {file_ext}: {e}")
            return None

    def _generate_pptx_thumbnail(self, file_path: str) -> Optional[bytes]:
        """
        Generate thumbnail from PowerPoint presentation
        Extracts the first slide's thumbnail if embedded, otherwise creates a placeholder

        Args:
            file_path: Path to PPTX file

        Returns:
            Thumbnail as bytes or None
        """
        try:
            import zipfile

            # PPTX files are ZIP archives - try to extract thumbnail
            with zipfile.ZipFile(file_path, 'r') as zip_ref:
                # PPTX contains thumbnail at docProps/thumbnail.jpeg
                thumbnail_paths = [
                    'docProps/thumbnail.jpeg',
                    'docProps/thumbnail.png',
                    '_rels/.rels'  # fallback
                ]

                for thumb_path in thumbnail_paths:
                    if thumb_path in zip_ref.namelist():
                        if 'thumbnail' in thumb_path:
                            with zip_ref.open(thumb_path) as thumb_file:
                                image_data = thumb_file.read()

                                # Resize to configured thumbnail size
                                image = Image.open(io.BytesIO(image_data))
                                image.thumbnail(
                                    (
                                        self.config.thumbnail.thumbnail_width,
                                        self.config.thumbnail.thumbnail_height
                                    ),
                                    Image.Resampling.LANCZOS
                                )

                                # Convert to RGB if necessary
                                if image.mode not in ('RGB', 'L'):
                                    image = image.convert('RGB')

                                # Save to bytes
                                buffer = io.BytesIO()
                                image.save(
                                    buffer,
                                    format=self.config.thumbnail.thumbnail_format,
                                    quality=self.config.thumbnail.thumbnail_quality
                                )

                                self.logger.debug("Generated thumbnail from PPTX embedded thumbnail")
                                return buffer.getvalue()

            self.logger.debug("No embedded thumbnail found in PPTX")
            return None

        except Exception as e:
            self.logger.warning(f"Failed to extract PPTX thumbnail: {e}")
            return None

    def _generate_office_thumbnail_via_libreoffice(self, file_path: str) -> Optional[bytes]:
        """
        Generate thumbnail using LibreOffice headless conversion
        Converts first page to PDF, then to image

        Args:
            file_path: Path to Office document

        Returns:
            Thumbnail as bytes or None
        """
        import subprocess
        import tempfile
        import os

        try:
            # Check if LibreOffice is available
            result = subprocess.run(
                ['which', 'libreoffice'],
                capture_output=True,
                text=True
            )

            if result.returncode != 0:
                self.logger.debug("LibreOffice not available for thumbnail generation")
                return None

            # Create temp directory for conversion
            with tempfile.TemporaryDirectory() as temp_dir:
                # Convert to PDF using LibreOffice headless
                cmd = [
                    'libreoffice',
                    '--headless',
                    '--convert-to', 'pdf',
                    '--outdir', temp_dir,
                    file_path
                ]

                result = subprocess.run(
                    cmd,
                    capture_output=True,
                    text=True,
                    timeout=60  # 60 second timeout
                )

                if result.returncode != 0:
                    self.logger.warning(f"LibreOffice conversion failed: {result.stderr}")
                    return None

                # Find the generated PDF
                base_name = Path(file_path).stem
                pdf_path = os.path.join(temp_dir, f"{base_name}.pdf")

                if not os.path.exists(pdf_path):
                    self.logger.warning("PDF output not found after LibreOffice conversion")
                    return None

                # Convert first page of PDF to image
                from pdf2image import convert_from_path

                images = convert_from_path(
                    pdf_path,
                    dpi=150,
                    first_page=1,
                    last_page=1
                )

                if not images:
                    return None

                image = images[0]

                # Resize to thumbnail
                image.thumbnail(
                    (
                        self.config.thumbnail.thumbnail_width,
                        self.config.thumbnail.thumbnail_height
                    ),
                    Image.Resampling.LANCZOS
                )

                # Convert to RGB if necessary
                if image.mode not in ('RGB', 'L'):
                    image = image.convert('RGB')

                # Save to bytes
                buffer = io.BytesIO()
                image.save(
                    buffer,
                    format=self.config.thumbnail.thumbnail_format,
                    quality=self.config.thumbnail.thumbnail_quality
                )

                self.logger.debug("Generated thumbnail via LibreOffice conversion")
                return buffer.getvalue()

        except subprocess.TimeoutExpired:
            self.logger.warning("LibreOffice conversion timed out")
            return None
        except Exception as e:
            self.logger.warning(f"LibreOffice thumbnail generation failed: {e}")
            return None
