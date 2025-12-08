"""
Pytest Configuration and Shared Fixtures
Provides common fixtures and setup for all tests
"""

import os
import sys
import tempfile
import shutil
from pathlib import Path
from typing import Generator, Dict, Any
from unittest.mock import Mock, MagicMock

import pytest
import boto3
from moto import mock_aws
from PIL import Image

# Add parent directory to path for imports
sys.path.insert(0, str(Path(__file__).parent.parent))

from config import Config, AWSConfig, ProcessingConfig, LoggingConfig


# ==========================================
# Session-level Fixtures
# ==========================================

@pytest.fixture(scope="session")
def test_data_dir() -> Path:
    """Get test data directory"""
    data_dir = Path(__file__).parent / "data"
    data_dir.mkdir(exist_ok=True)
    return data_dir


@pytest.fixture(scope="session")
def temp_dir() -> Generator[Path, None, None]:
    """Create temporary directory for tests"""
    tmp_dir = Path(tempfile.mkdtemp(prefix="test_worker_"))
    yield tmp_dir
    # Cleanup
    shutil.rmtree(tmp_dir, ignore_errors=True)


# ==========================================
# Configuration Fixtures
# ==========================================

@pytest.fixture
def test_config() -> Config:
    """Create test configuration"""
    config = Config()

    # AWS configuration
    config.aws = AWSConfig(
        region="us-east-1",
        s3_bucket="test-bucket",
        sqs_queue_url="https://sqs.us-east-1.amazonaws.com/123456789/test-queue",
        opensearch_endpoint="https://test-opensearch.us-east-1.es.amazonaws.com",
        opensearch_index="test-index",
    )

    # Processing configuration
    config.processing = ProcessingConfig(
        max_file_size_mb=100,
        temp_dir="/tmp/test_worker",
        max_workers=2,
        batch_size=5,
    )

    # Logging configuration
    config.logging = LoggingConfig(
        log_level="DEBUG",
        log_file="/tmp/test_worker.log",
    )

    return config


# ==========================================
# AWS Mock Fixtures
# ==========================================

@pytest.fixture
def aws_credentials():
    """Mock AWS credentials for testing"""
    os.environ['AWS_ACCESS_KEY_ID'] = 'testing'
    os.environ['AWS_SECRET_ACCESS_KEY'] = 'testing'
    os.environ['AWS_SECURITY_TOKEN'] = 'testing'
    os.environ['AWS_SESSION_TOKEN'] = 'testing'
    os.environ['AWS_DEFAULT_REGION'] = 'us-east-1'


@pytest.fixture
def s3_client(aws_credentials):
    """Create mock S3 client"""
    with mock_aws():
        s3 = boto3.client('s3', region_name='us-east-1')
        # Create test bucket
        s3.create_bucket(Bucket='test-bucket')
        yield s3


@pytest.fixture
def sqs_client(aws_credentials):
    """Create mock SQS client"""
    with mock_aws():
        sqs = boto3.client('sqs', region_name='us-east-1')
        # Create test queue
        response = sqs.create_queue(QueueName='test-queue')
        queue_url = response['QueueUrl']
        yield sqs, queue_url


@pytest.fixture
def mock_opensearch():
    """Mock OpenSearch client"""
    mock_client = MagicMock()
    mock_client.ping.return_value = True
    mock_client.indices.exists.return_value = False
    mock_client.indices.create.return_value = {'acknowledged': True}
    mock_client.index.return_value = {
        '_index': 'test-index',
        '_id': 'test-id',
        'result': 'created'
    }
    return mock_client


# ==========================================
# Test File Fixtures
# ==========================================

@pytest.fixture
def sample_pdf(temp_dir: Path) -> Path:
    """Create sample PDF file for testing"""
    from pypdf import PdfWriter, PdfReader
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    import io

    # Create PDF with reportlab
    buffer = io.BytesIO()
    c = canvas.Canvas(buffer, pagesize=letter)
    c.drawString(100, 750, "Test PDF Document")
    c.drawString(100, 730, "This is a sample PDF for testing.")
    c.drawString(100, 710, "It contains multiple lines of text.")
    c.showPage()
    c.save()

    # Write to file
    pdf_path = temp_dir / "sample.pdf"
    buffer.seek(0)
    pdf_path.write_bytes(buffer.getvalue())

    return pdf_path


@pytest.fixture
def sample_image(temp_dir: Path) -> Path:
    """Create sample image file for testing"""
    img_path = temp_dir / "sample.jpg"

    # Create simple image
    img = Image.new('RGB', (800, 600), color='white')
    img.save(img_path, 'JPEG')

    return img_path


@pytest.fixture
def sample_text_image(temp_dir: Path) -> Path:
    """Create image with text for OCR testing"""
    from PIL import ImageDraw, ImageFont

    img_path = temp_dir / "text_image.png"

    # Create image with text
    img = Image.new('RGB', (800, 200), color='white')
    draw = ImageDraw.Draw(img)

    # Draw text
    text = "Sample text for OCR testing"
    draw.text((50, 80), text, fill='black')

    img.save(img_path, 'PNG')

    return img_path


@pytest.fixture
def sample_docx(temp_dir: Path) -> Path:
    """Create sample Word document for testing"""
    from docx import Document

    doc_path = temp_dir / "sample.docx"

    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a sample Word document.')
    doc.add_paragraph('It contains multiple paragraphs.')
    doc.add_heading('Section 1', level=1)
    doc.add_paragraph('Content of section 1.')

    doc.save(doc_path)

    return doc_path


@pytest.fixture
def sample_xlsx(temp_dir: Path) -> Path:
    """Create sample Excel file for testing"""
    from openpyxl import Workbook

    xlsx_path = temp_dir / "sample.xlsx"

    wb = Workbook()
    ws = wb.active
    ws.title = "Test Sheet"

    # Add data
    ws['A1'] = 'Name'
    ws['B1'] = 'Value'
    ws['A2'] = 'Item 1'
    ws['B2'] = 100
    ws['A3'] = 'Item 2'
    ws['B3'] = 200

    wb.save(xlsx_path)

    return xlsx_path


@pytest.fixture
def sample_pptx(temp_dir: Path) -> Path:
    """Create sample PowerPoint file for testing"""
    from pptx import Presentation
    from pptx.util import Inches

    pptx_path = temp_dir / "sample.pptx"

    prs = Presentation()

    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Test Presentation"
    subtitle.text = "Sample PowerPoint for testing"

    # Add content slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = 'Content'
    tf = body_shape.text_frame
    tf.text = 'First bullet point'

    prs.save(pptx_path)

    return pptx_path


# ==========================================
# SQS Message Fixtures
# ==========================================

@pytest.fixture
def s3_event_message() -> Dict[str, Any]:
    """Create sample S3 event message"""
    return {
        "Records": [
            {
                "eventVersion": "2.1",
                "eventSource": "aws:s3",
                "awsRegion": "us-east-1",
                "eventTime": "2024-01-01T00:00:00.000Z",
                "eventName": "ObjectCreated:Put",
                "s3": {
                    "bucket": {
                        "name": "test-bucket",
                        "arn": "arn:aws:s3:::test-bucket"
                    },
                    "object": {
                        "key": "test-files/sample.pdf",
                        "size": 1024,
                        "eTag": "abc123"
                    }
                }
            }
        ]
    }


@pytest.fixture
def custom_sqs_message() -> Dict[str, Any]:
    """Create custom SQS message format"""
    return {
        "bucket": "test-bucket",
        "key": "test-files/sample.pdf",
        "timestamp": "2024-01-01T00:00:00.000Z",
        "metadata": {
            "source": "test",
            "priority": "normal"
        }
    }


# ==========================================
# Processing Result Fixtures
# ==========================================

@pytest.fixture
def mock_processing_result():
    """Create mock processing result"""
    from processors.base_processor import ProcessingResult

    return ProcessingResult(
        success=True,
        file_path="/tmp/test.pdf",
        file_name="test.pdf",
        file_size=1024,
        file_type=".pdf",
        mime_type="application/pdf",
        extracted_text="Sample extracted text",
        page_count=1,
        word_count=3,
        char_count=21,
        processing_time_seconds=0.5,
        processor_name="TestProcessor",
    )


# ==========================================
# Helper Functions
# ==========================================

@pytest.fixture
def create_test_file():
    """Factory fixture to create test files"""
    def _create_file(path: Path, content: str = "test content") -> Path:
        path.write_text(content)
        return path
    return _create_file


@pytest.fixture
def upload_to_s3():
    """Helper to upload files to mock S3"""
    def _upload(s3_client, bucket: str, key: str, file_path: Path):
        s3_client.upload_file(str(file_path), bucket, key)
        return f"s3://{bucket}/{key}"
    return _upload


@pytest.fixture
def send_sqs_message():
    """Helper to send messages to mock SQS"""
    def _send(sqs_client, queue_url: str, message_body: Dict[str, Any]):
        import json
        response = sqs_client.send_message(
            QueueUrl=queue_url,
            MessageBody=json.dumps(message_body)
        )
        return response['MessageId']
    return _send


# ==========================================
# Performance Testing Fixtures
# ==========================================

@pytest.fixture
def performance_config():
    """Configuration for performance tests"""
    return {
        'max_processing_time': 5.0,  # seconds
        'max_memory_mb': 500,
        'min_throughput': 10,  # files per second
    }


# ==========================================
# Error Injection Fixtures
# ==========================================

@pytest.fixture
def inject_s3_error():
    """Inject S3 errors for fault tolerance testing"""
    from botocore.exceptions import ClientError

    def _inject_error(error_code: str = "NoSuchKey"):
        error = ClientError(
            {
                'Error': {
                    'Code': error_code,
                    'Message': 'Test error'
                }
            },
            'GetObject'
        )
        return error
    return _inject_error


# ==========================================
# Cleanup
# ==========================================

@pytest.fixture(autouse=True)
def cleanup_temp_files(temp_dir: Path):
    """Cleanup temporary files after each test"""
    yield
    # Clean up temp directory contents
    for item in temp_dir.glob('*'):
        if item.is_file():
            item.unlink()
        elif item.is_dir():
            shutil.rmtree(item)


# ==========================================
# Custom Assertions
# ==========================================

class ProcessingResultAssertion:
    """Custom assertions for ProcessingResult"""

    @staticmethod
    def assert_success(result):
        """Assert processing was successful"""
        assert result.success, f"Processing failed: {result.error_message}"
        assert result.error_message is None
        assert result.extracted_text
        assert result.char_count > 0

    @staticmethod
    def assert_failure(result, error_pattern: str = None):
        """Assert processing failed with expected error"""
        assert not result.success
        assert result.error_message is not None
        if error_pattern:
            assert error_pattern in result.error_message


@pytest.fixture
def assert_processing():
    """Fixture for processing result assertions"""
    return ProcessingResultAssertion()
